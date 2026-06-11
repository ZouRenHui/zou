#!/usr/bin/env python3
"""从 Word / TXT / PDF / PPT 等文档中提取纯文本。"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

SUPPORTED_SUFFIXES = {".txt", ".md", ".docx", ".pdf", ".pptx", ".doc", ".ppt", ".rtf"}


def extract_text(path: Path) -> str:
    """读取文档正文，返回纯文本。"""
    if not path.is_file():
        raise FileNotFoundError(f"找不到文件: {path}")

    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        raise ValueError(f"不支持的文件类型: {suffix}")

    if suffix in {".txt", ".md"}:
        return _read_text_file(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix in {".doc", ".ppt", ".rtf"}:
        return _read_legacy_office(path)
    raise ValueError(f"不支持的文件类型: {suffix}")


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("未安装 python-docx，请运行: pip install python-docx") from exc

    doc = Document(str(path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("未安装 pypdf，请运行: pip install pypdf") from exc

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("未安装 python-pptx，请运行: pip install python-pptx") from exc

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                slide_parts.append(text)
        if slide_parts:
            parts.append(f"【幻灯片 {slide_index}】\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def _read_legacy_office(path: Path) -> str:
    if sys.platform == "darwin":
        return _read_with_textutil(path)
    raise RuntimeError(
        f"旧版 Office 格式 ({path.suffix}) 建议先另存为 .docx/.pptx，"
        "或在 macOS 上使用系统 textutil 转换。"
    )


def _read_with_textutil(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".ppt":
        raise RuntimeError("macOS textutil 不支持 .ppt，请先转换为 .pptx。")

    result = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", str(path)],
        capture_output=True,
        text=True,
        check=False,
    )
    if result.returncode != 0:
        stderr = (result.stderr or "").strip() or "textutil 转换失败"
        raise RuntimeError(stderr)
    return result.stdout.strip()
